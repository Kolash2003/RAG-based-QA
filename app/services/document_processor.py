from pathlib import Path
from typing import List, Dict
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
import pandas as pd
from app.utils.logger import logger
from app.config import settings


class DocumentProcessor:
    """Handles document parsing and text extraction."""
    
    SUPPORTED_FORMATS = {
        '.pdf': 'pdf',
        '.txt': 'text',
        '.docx': 'docx',
        '.pptx': 'pptx',
        '.xlsx': 'xlsx',
        '.csv': 'csv'
    }
    
    def __init__(self):
        self.chunk_size = settings.chunk_size
        self.chunk_overlap = settings.chunk_overlap
    
    def is_supported(self, filename: str) -> bool:
        """Check if file format is supported."""
        return Path(filename).suffix.lower() in self.SUPPORTED_FORMATS
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from document based on file type."""
        path = Path(file_path)
        suffix = path.suffix.lower()
        
        try:
            if suffix == '.pdf':
                return self._extract_from_pdf(file_path)
            elif suffix == '.txt':
                return self._extract_from_txt(file_path)
            elif suffix == '.docx':
                return self._extract_from_docx(file_path)
            elif suffix == '.pptx':
                return self._extract_from_pptx(file_path)
            elif suffix == '.xlsx':
                return self._extract_from_xlsx(file_path)
            elif suffix == '.csv':
                return self._extract_from_csv(file_path)
            else:
                raise ValueError(f"Unsupported file format: {suffix}")
        except Exception as e:
            logger.error("text_extraction_failed", file_path=file_path, error=str(e))
            raise
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF."""
        text = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text())
        return "\n".join(text)
    
    def _extract_from_txt(self, file_path: str) -> str:
        """Extract text from TXT file."""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX."""
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    
    def _extract_from_xlsx(self, file_path: str) -> str:
        """Extract text from XLSX."""
        wb = openpyxl.load_workbook(file_path)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text.append(" | ".join([str(cell) for cell in row if cell]))
        return "\n".join(text)
    
    def _extract_from_csv(self, file_path: str) -> str:
        """Extract text from CSV."""
        df = pd.read_csv(file_path)
        return df.to_string()
    
    def chunk_text(self, text: str, metadata: Dict = None) -> List[Dict]:
        """Split text into chunks with overlap."""
        if not text:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + self.chunk_size
            chunk_text = text[start:end]

            if end < text_length:
                last_period = chunk_text.rfind('.')
                last_newline = chunk_text.rfind('\n')
                last_space = chunk_text.rfind(' ')
                
                break_point = max(last_period, last_newline, last_space)
                if break_point > self.chunk_size * 0.8:  # At least 80% of chunk size
                    end = start + break_point + 1
                    chunk_text = text[start:end]
            
            chunk = {
                "text": chunk_text.strip(),
                "metadata": metadata or {},
                "chunk_index": len(chunks)
            }
            chunks.append(chunk)
            
            start = end - self.chunk_overlap
        
        logger.info("text_chunked", num_chunks=len(chunks), text_length=text_length)
        return chunks


document_processor = DocumentProcessor()